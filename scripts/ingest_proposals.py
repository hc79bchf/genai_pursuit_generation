import os
import glob
import logging
import argparse
from typing import List, Dict
import chromadb
from chromadb.utils import embedding_functions
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from a .docx file."""
    try:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    except Exception as e:
        logger.exception(f"Error reading {file_path}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a .pptx file."""
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
        return "\n".join(text_runs)
    except Exception as e:
        logger.exception(f"Error reading {file_path}")
        return ""

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 100) -> List[str]:
    """
    Split text into chunks with overlap, respecting sentence boundaries where possible.
    This is a simple improvement over character slicing.
    """
    if not text:
        return []
    
    # Simple recursive splitting by paragraphs first, then characters if needed
    # For now, let's improve the sliding window to not cut words in half if possible
    # A more robust solution would use langchain's RecursiveCharacterTextSplitter
    
    chunks = []
    start = 0
    text_len = len(text)
    
    while start < text_len:
        end = start + chunk_size
        
        # If we are not at the end of text, try to find the last space within the chunk
        if end < text_len:
            # Look for the last space in the chunk to avoid splitting words
            last_space = text.rfind(' ', start, end)
            if last_space != -1:
                end = last_space
        
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        
        # Move start forward, subtracting overlap
        start = end - overlap
        
        # Ensure we always move forward to avoid infinite loops
        if start >= end:
            start = end
            
    return chunks

def main():
    parser = argparse.ArgumentParser(description="Ingest proposals into ChromaDB")
    parser.add_argument("--data-dir", default=os.path.join(os.path.dirname(__file__), "../Data/PriorProposal"), help="Directory containing proposal files")
    parser.add_argument("--chroma-host", default="localhost", help="ChromaDB host")
    parser.add_argument("--chroma-port", type=int, default=8001, help="ChromaDB port")
    parser.add_argument("--collection-name", default="prior_proposals", help="ChromaDB collection name")
    args = parser.parse_args()

    openai_api_key = os.getenv("OPENAI_API_KEY")
    if not openai_api_key:
        logger.error("OPENAI_API_KEY not found in environment variables")
        return

    logger.info(f"Connecting to ChromaDB at {args.chroma_host}:{args.chroma_port}...")
    try:
        client = chromadb.HttpClient(host=args.chroma_host, port=args.chroma_port)
        logger.info("Successfully connected to ChromaDB.")
    except Exception as e:
        logger.error(f"Failed to connect to ChromaDB HTTP client: {e}")
        logger.info("Falling back to local PersistentClient at ./chroma_db")
        client = chromadb.PersistentClient(path="./chroma_db")

    # Initialize OpenAI embedding function
    openai_ef = embedding_functions.OpenAIEmbeddingFunction(
        api_key=openai_api_key,
        model_name="text-embedding-3-small"
    )

    # Get or create collection
    try:
        collection = client.get_or_create_collection(
            name=args.collection_name,
            embedding_function=openai_ef
        )
        logger.info(f"Collection '{args.collection_name}' ready.")
    except Exception as e:
        logger.exception(f"Error getting/creating collection")
        return

    # Process files
    if not os.path.exists(args.data_dir):
        logger.error(f"Data directory not found: {args.data_dir}")
        return

    files = glob.glob(os.path.join(args.data_dir, "*"))
    logger.info(f"Found {len(files)} files in {args.data_dir}")

    total_chunks = 0
    
    for file_path in files:
        filename = os.path.basename(file_path)
        logger.info(f"Processing {filename}...")
        
        if filename.endswith(".docx"):
            text = extract_text_from_docx(file_path)
        elif filename.endswith(".pptx"):
            text = extract_text_from_pptx(file_path)
        else:
            logger.warning(f"Skipping unsupported file type: {filename}")
            continue
            
        if not text:
            logger.warning(f"No text extracted from {filename}")
            continue
            
        chunks = chunk_text(text)
        logger.info(f"  - Generated {len(chunks)} chunks")
        
        if not chunks:
            continue

        # Prepare data for ChromaDB
        ids = [f"{filename}_chunk_{i}" for i in range(len(chunks))]
        metadatas = [{"source": filename, "chunk_index": i} for i in range(len(chunks))]
        
        # Add to collection
        try:
            collection.add(
                documents=chunks,
                metadatas=metadatas,
                ids=ids
            )
            total_chunks += len(chunks)
            logger.info(f"  - Indexed {len(chunks)} chunks")
        except Exception as e:
            logger.exception(f"  - Error indexing chunks for {filename}")

    logger.info(f"Ingestion complete! Total chunks indexed: {total_chunks}")

if __name__ == "__main__":
    main()
